"""
preview.py — PPTX 布局坐标可视化工具
用法：
  python preview.py <文件路径> --slide N
  python preview.py <文件路径> --slide N --out preview.png
  python preview.py <文件路径> --slide N --show          # 直接弹窗显示

依赖：pip install python-pptx matplotlib
"""

import sys
import io
import argparse
from pathlib import Path

# ⚠️ inspect.py 与标准库 inspect 同名，scripts/ 目录在 sys.path[0] 时干扰 pptx import
# preview.py 与 inspect.py 同目录，必须同样移除该目录
import sys as _sys
import os as _os
_this_dir = _os.path.normcase(_os.path.abspath(_os.path.dirname(__file__)))
_sys.path = [p for p in _sys.path if _os.path.normcase(_os.path.abspath(p)) != _this_dir]
del _this_dir

# ── 依赖 import 必须在 stdout 重包装之前 ──────────────────────────
def _err(msg):
    """import 失败时的安全输出（stdout 尚未重包装）"""
    raw = (msg + "\n").encode('utf-8')
    sys.stdout.buffer.write(raw) if hasattr(sys.stdout, 'buffer') else print(msg)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:
    _err(f"❌ 缺少依赖：python-pptx 未安装\n   请执行：{sys.executable} -m pip install python-pptx")

try:
    import matplotlib
    import matplotlib.pyplot as plt
    import matplotlib.patches as patches
    from matplotlib.patches import FancyBboxPatch
    # ── 中文字体自动配置（Windows 优先 Microsoft YaHei，macOS 用 PingFang SC）──
    import matplotlib.font_manager as _fm
    _CN_FONTS = ['Microsoft YaHei', 'SimHei', 'PingFang SC', 'Heiti TC', 'Arial Unicode MS']
    _available = {f.name for f in _fm.fontManager.ttflist}
    _chosen = next((f for f in _CN_FONTS if f in _available), None)
    if _chosen:
        matplotlib.rcParams['font.family'] = _chosen
    matplotlib.rcParams['axes.unicode_minus'] = False  # 防止负号乱码
except ImportError:
    _err(f"❌ 缺少依赖：matplotlib 未安装\n   请执行：{sys.executable} -m pip install matplotlib")

# 依赖 import 成功后再重包装 stdout
if __name__ == "__main__" and hasattr(sys.stdout, 'buffer'):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
if __name__ == "__main__" and hasattr(sys.stderr, 'buffer'):
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
# ──────────────────────────────────────────────────────────────────


def emu_to_inch(emu):
    return emu / 914400 if emu is not None else 0


def get_fill_hex(shape):
    """尝试获取 shape 填充色，返回 matplotlib 颜色字符串"""
    try:
        c = shape.fill.fore_color
        if c and c.type:
            rgb = c.rgb
            return f"#{rgb}"
    except Exception:
        pass
    return None


# 固定边界元素的识别（同 inspect.py 逻辑）
def is_boundary_element(shape, slide_w, slide_h):
    h = emu_to_inch(shape.height)
    w = emu_to_inch(shape.width)
    l = emu_to_inch(shape.left)
    t = emu_to_inch(shape.top)
    if h >= slide_h * 0.8:
        return True
    area = w * h
    corner_margin = 1.0
    in_corner = (
        (l < corner_margin or l + w > slide_w - corner_margin) and
        (t < corner_margin or t + h > slide_h - corner_margin)
    )
    if area < 0.5 and in_corner and not shape.has_text_frame:
        return True
    return False


def draw_slide(prs, slide_idx, output_path=None, show=False):
    slide   = prs.slides[slide_idx]
    slide_w = emu_to_inch(prs.slide_width)
    slide_h = emu_to_inch(prs.slide_height)

    # 画布：比例与幻灯片一致，宽度固定 14 英寸
    fig_w   = 14
    fig_h   = fig_w * slide_h / slide_w
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.set_xlim(0, slide_w)
    ax.set_ylim(0, slide_h)
    ax.invert_yaxis()   # PPT 坐标系：y 向下为正
    ax.set_aspect('equal')
    ax.set_facecolor('#F5F5F0')
    ax.set_title(f"第 {slide_idx+1} 页 — 布局预览  ({slide_w}\" × {slide_h}\")",
                 fontsize=11, pad=8)
    ax.set_xlabel("← 横向（英寸）→", fontsize=8)
    ax.set_ylabel("← 纵向（英寸）→", fontsize=8)

    # 画网格
    import numpy as np
    for x in np.arange(0, slide_w + 0.5, 0.5):
        ax.axvline(x, color='#CCCCCC', linewidth=0.3, zorder=0)
    for y in np.arange(0, slide_h + 0.5, 0.5):
        ax.axhline(y, color='#CCCCCC', linewidth=0.3, zorder=0)

    # 颜色方案
    TYPE_COLORS = {
        "TEXT_BOX":   ('#D6EAF8', '#2E86C1'),   # 浅蓝/蓝
        "AUTO_SHAPE": ('#D5F5E3', '#1E8449'),    # 浅绿/绿
        "PICTURE":    ('#FDEBD0', '#CA6F1E'),    # 浅橙/橙
        "TABLE":      ('#F9EBEA', '#C0392B'),    # 浅红/红
        "CHART":      ('#F4ECF7', '#7D3C98'),    # 浅紫/紫
        "PLACEHOLDER":('#FDFEFE', '#717D7E'),    # 浅灰/灰
        "GROUP":      ('#FEFEFE', '#AAB7B8'),
        "DEFAULT":    ('#FDFDFD', '#808B96'),
    }
    BOUNDARY_COLOR = ('#FADBD8', '#E74C3C')      # 红色：固定边界元素

    for idx, shape in enumerate(slide.shapes):
        l = emu_to_inch(shape.left)
        t = emu_to_inch(shape.top)
        w = emu_to_inch(shape.width)
        h = emu_to_inch(shape.height)

        if w <= 0 or h <= 0:
            continue

        is_boundary = is_boundary_element(shape, slide_w, slide_h)
        shape_type  = {
            1: "AUTO_SHAPE", 4: "CHART", 7: "GROUP", 15: "PICTURE",
            16: "PLACEHOLDER", 18: "TABLE", 19: "TEXT_BOX",
        }.get(shape.shape_type, "DEFAULT")

        fc, ec = BOUNDARY_COLOR if is_boundary else TYPE_COLORS.get(shape_type, TYPE_COLORS["DEFAULT"])

        # 优先用 shape 自身填充色（半透明叠加）
        fill_hex = get_fill_hex(shape)
        if fill_hex and not is_boundary:
            fc = fill_hex

        rect = FancyBboxPatch(
            (l, t), w, h,
            boxstyle="round,pad=0.01",
            facecolor=fc, edgecolor=ec,
            linewidth=1.2 if is_boundary else 0.8,
            alpha=0.75,
            zorder=2,
        )
        ax.add_patch(rect)

        # 标签：shape_id + 前12字文本
        label_parts = [f"id={shape.shape_id}"]
        if shape.has_text_frame:
            txt = shape.text_frame.text[:12].replace('\n', '/')
            if txt.strip():
                label_parts.append(txt)
        if is_boundary:
            label_parts.append("[边界]")

        label    = '\n'.join(label_parts)
        font_sz  = max(4.5, min(7, w * 6))  # 字号随框宽缩放
        text_x   = l + w / 2
        text_y   = t + h / 2

        ax.text(text_x, text_y, label,
                ha='center', va='center',
                fontsize=font_sz,
                color='#1A1A1A',
                wrap=True,
                zorder=3,
                bbox=dict(boxstyle='round,pad=0.1', fc='white', alpha=0.6, ec='none'))

    # 图例
    legend_items = [
        patches.Patch(facecolor='#D6EAF8', edgecolor='#2E86C1', label='TEXT_BOX'),
        patches.Patch(facecolor='#D5F5E3', edgecolor='#1E8449', label='AUTO_SHAPE'),
        patches.Patch(facecolor='#FDEBD0', edgecolor='#CA6F1E', label='PICTURE'),
        patches.Patch(facecolor='#F9EBEA', edgecolor='#C0392B', label='TABLE'),
        patches.Patch(facecolor='#F4ECF7', edgecolor='#7D3C98', label='CHART'),
        patches.Patch(facecolor='#FADBD8', edgecolor='#E74C3C', label='固定边界元素 [边界]'),
    ]
    ax.legend(handles=legend_items, loc='upper right',
              fontsize=7, framealpha=0.85, bbox_to_anchor=(1.0, -0.04),
              ncol=3)

    plt.tight_layout()

    if output_path:
        plt.savefig(output_path, dpi=150, bbox_inches='tight')
        print(f"✅ 预览图已保存：{output_path}")
    if show:
        # backend 已在 main() 中通过 matplotlib.use('TkAgg'/'Qt5Agg') 设置
        # 此处直接调用 plt.show()，不重复 use()（use() 在 figure 创建后调用无效）
        plt.show()
    if not output_path and not show:
        # 默认保存到与 pptx 同目录
        return fig

    plt.close(fig)


def main():
    parser = argparse.ArgumentParser(description="PPTX 布局坐标可视化工具")
    parser.add_argument("file",              help="PPTX 文件路径")
    parser.add_argument("--slide",  type=int, default=1, help="分析第N页（1-based，默认第1页）")
    parser.add_argument("--out",             help="输出 PNG 路径（默认与 PPTX 同目录同名）")
    parser.add_argument("--show",   action="store_true", help="直接弹窗显示（需要图形界面）")
    parser.add_argument("--all",    action="store_true", help="输出所有页（每页一个 PNG）")
    args = parser.parse_args()

    # ── backend 选择：--show 时用交互式，否则用 Agg（无需图形界面）──
    if args.show:
        try:
            matplotlib.use('TkAgg')
        except Exception:
            try:
                matplotlib.use('Qt5Agg')
            except Exception:
                print("⚠️  无法启动图形界面，改为保存 PNG 文件")
                args.show = False
    else:
        matplotlib.use('Agg')

    pptx_path = Path(args.file)
    if not pptx_path.exists():
        print(f"❌ 文件不存在：{pptx_path}")
        sys.exit(1)

    prs = Presentation(str(pptx_path))

    if args.all:
        for i in range(len(prs.slides)):
            out = pptx_path.parent / f"{pptx_path.stem}_preview_slide{i+1}.png"
            draw_slide(prs, i, output_path=str(out))
    else:
        slide_idx = args.slide - 1
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            print(f"❌ 页码超出范围（共 {len(prs.slides)} 页）")
            sys.exit(1)
        out = args.out or str(pptx_path.parent / f"{pptx_path.stem}_preview_slide{args.slide}.png")
        draw_slide(prs, slide_idx, output_path=out, show=args.show)


if __name__ == "__main__":
    main()
