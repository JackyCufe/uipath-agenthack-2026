"""
inspect.py — PPTX 结构完整分析工具
用法：
  python inspect.py <文件路径>               # 分析所有页
  python inspect.py <文件路径> --slide N     # 只分析第N页
  python inspect.py <文件路径> --layouts     # 额外列出所有可用布局
  python inspect.py <文件路径> --json        # 输出 JSON 格式（供脚本使用）
  python inspect.py <文件路径> --relations   # 输出 shape 空间关联分析
  python inspect.py <文件路径> --overflow    # 输出文字溢出风险评估
"""

# ═══════════════════════════════════════════════════════════════════
# ⚠️  同名保护（必须是文件第一段可执行代码）
#
# 本文件名 inspect.py 与 Python 标准库 inspect 同名。
# 当 scripts/ 目录在 sys.path 中时（直接运行本脚本时 Python 会把脚本所在目录
# 加入 sys.path[0]），lxml / pptx 等库内部的 `import inspect` 会命中本文件，
# 导致 lxml C 扩展初始化时找不到 getfullargspec/getargspec 等标准函数而崩溃。
#
# 修复策略（两层）：
#   1. 若本文件不是以 __main__ 直接运行（即被误当标准库加载），
#      立即把 sys.modules['inspect'] 替换成真正的标准库，然后让调用方透明地
#      拿到正确的 inspect —— 本文件自身不再继续初始化。
#   2. 同时把 scripts/ 从 sys.path 中移除，防止后续 import 继续命中本文件。
# ═══════════════════════════════════════════════════════════════════
import sys as _sys, os as _os

if __name__ != "__main__":
    # 被误当标准库 inspect import：
    # 把 scripts/ 从 sys.path 移除，将真正标准库的所有符号注入本模块命名空间，
    # lxml/pptx 等调用方透明拿到正确的 inspect（有 getfullargspec 等函数）。
    _sys.modules.pop("inspect", None)
    _this_dir = _os.path.normcase(_os.path.abspath(_os.path.dirname(__file__) or "."))
    _sys.path = [p for p in _sys.path if _os.path.normcase(_os.path.abspath(p)) != _this_dir]
    import importlib as _il
    _real = _il.import_module("inspect")
    _sys.modules["inspect"] = _real
    globals().update({k: v for k, v in vars(_real).items() if not k.startswith("__")})
    del _il, _real, _this_dir
    # 不再执行后续 PPTX 工具代码

# ── 以下仅在直接运行（python inspect.py ...）时执行 ────────────────
# 移除 scripts/ 目录，防止后续 from pptx import ... 再次命中本文件
_this_dir = _os.path.normcase(_os.path.abspath(_os.path.dirname(__file__) or "."))
_sys.path = [p for p in _sys.path if _os.path.normcase(_os.path.abspath(p)) != _this_dir]
del _this_dir

import sys
import io
import json
import argparse
import math
import os
from pathlib import Path

# ── 依赖检测（修复版）──────────────────────────────────────────────
# 修复要点：
# 1. import 放在 stdout 重包装之前，确保失败时仍能正常输出错误信息
# 2. import 失败时输出当前解释器路径，方便在多 Python 环境下定位问题
# 3. stdout 重包装仅在 buffer 属性存在时执行，避免管道重定向场景下报错
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:
    # 仅在「直接运行本脚本」时把缺依赖当致命错误退出。
    # 若本文件被其他代码以标准库名 `inspect` 误 import（scripts 目录混入 sys.path 时，
    # html5lib / webencodings 等内部的 `import inspect` 会命中这里），绝不能 sys.exit()
    # 杀掉整个进程——置 None 延迟到真正使用时再报错。
    Presentation = None  # type: ignore[assignment]
    RGBColor = None      # type: ignore[assignment]
    if __name__ == "__main__":
        interpreter = sys.executable
        # 此时 stdout 尚未重包装，直接写 buffer 以保证输出不乱码
        msg = (
            f"❌ 缺少依赖：python-pptx 未安装\n"
            f"   当前解释器：{interpreter}\n"
            f"   请执行：{interpreter} -m pip install python-pptx\n"
        )
        sys.stdout.buffer.write(msg.encode('utf-8')) if hasattr(sys.stdout, 'buffer') else print(msg)
        sys.exit(1)

# 依赖 import 成功后再重包装 stdout（保证中文路径正常输出）
if __name__ == "__main__" and hasattr(sys.stdout, 'buffer'):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
if __name__ == "__main__" and hasattr(sys.stderr, 'buffer'):
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
# ──────────────────────────────────────────────────────────────────

SHAPE_TYPE_NAMES = {
    1: "AUTO_SHAPE", 3: "CANVAS", 4: "CHART", 6: "FREEFORM",
    7: "GROUP", 11: "LINE", 13: "LINKED_PICTURE", 14: "OLE_OBJECT",
    15: "PICTURE", 16: "PLACEHOLDER", 18: "TABLE", 19: "TEXT_BOX",
}

# 中文字符宽度系数（相对于字号，单位 pt）
# 中文字符约等于 1.0 × 字号宽，英文/数字约 0.55 × 字号宽
CN_CHAR_WIDTH_RATIO  = 1.0
EN_CHAR_WIDTH_RATIO  = 0.55
PT_TO_INCH           = 1 / 72.0


def emu_to_inch(emu):
    return round(emu / 914400, 3) if emu is not None else None


def get_color_str(color_obj):
    try:
        if color_obj and color_obj.type:
            return f"#{color_obj.rgb}"
    except Exception:
        pass
    return None


def estimate_text_width_inch(text, font_size_pt):
    """粗略估算一行文字的渲染宽度（英寸）"""
    width_pt = 0.0
    for ch in text:
        if ord(ch) > 127:
            width_pt += font_size_pt * CN_CHAR_WIDTH_RATIO
        else:
            width_pt += font_size_pt * EN_CHAR_WIDTH_RATIO
    return width_pt * PT_TO_INCH


def _resolve_font_size_pt(run, para, shape) -> float:
    """
    解析 run 的字号（pt），按继承链向上查找：run → para → shape → 兜底 12pt。
    run.font.size / para.font.size / shape.text_frame 级别的默认字号均为 EMU * 12700。
    """
    # 1. run 自身有字号
    if run.font.size:
        return round(run.font.size / 12700, 1)
    # 2. 段落级默认字号
    try:
        if para.font.size:
            return round(para.font.size / 12700, 1)
    except Exception:
        pass
    # 3. shape 的 text_frame 默认字号（通过 XML 读取）
    try:
        from pptx.oxml.ns import qn as _qn
        txBody = shape.text_frame._txBody
        lstStyle = txBody.find(_qn('a:lstStyle'))
        if lstStyle is not None:
            lvl1 = lstStyle.find(_qn('a:lvl1pPr'))
            if lvl1 is not None:
                defRPr = lvl1.find(_qn('a:defRPr'))
                if defRPr is not None:
                    sz = defRPr.get('sz')
                    if sz:
                        return round(int(sz) / 100, 1)
    except Exception:
        pass
    # 4. 兜底：12pt（PowerPoint 正文默认值，比原来的 10pt 更接近实际）
    return 12.0


def assess_overflow_risk(shape):
    """
    评估文本框溢出风险。
    返回 dict: {risk: 'low'/'medium'/'high', reason: str, details: [...]}
    """
    if not shape.has_text_frame:
        return None

    tf       = shape.text_frame
    box_w    = emu_to_inch(shape.width) or 0
    box_h    = emu_to_inch(shape.height) or 0
    details  = []
    max_risk = 'low'

    for para in tf.paragraphs:
        if not para.text.strip():
            continue
        for run in para.runs:
            if not run.text.strip():
                continue
            font_pt = _resolve_font_size_pt(run, para, shape)
            est_w   = estimate_text_width_inch(run.text, font_pt)
            ratio   = est_w / box_w if box_w > 0 else 999

            if ratio > 1.2:
                risk = 'high'
            elif ratio > 0.9:
                risk = 'medium'
            else:
                risk = 'low'

            if risk != 'low':
                details.append({
                    "text":       run.text[:30],
                    "font_pt":    font_pt,
                    "est_width":  round(est_w, 3),
                    "box_width":  box_w,
                    "ratio":      round(ratio, 2),
                    "risk":       risk,
                })
                if risk == 'high':
                    max_risk = 'high'
                elif risk == 'medium' and max_risk == 'low':
                    max_risk = 'medium'

    return {"risk": max_risk, "details": details} if details else {"risk": "low", "details": []}


def inspect_shape(shape, idx):
    info = {
        "index": idx,
        "shape_id": shape.shape_id,
        "name": shape.name,
        "type": SHAPE_TYPE_NAMES.get(shape.shape_type, str(shape.shape_type)),
        "position": {
            "left":   emu_to_inch(shape.left),
            "top":    emu_to_inch(shape.top),
            "width":  emu_to_inch(shape.width),
            "height": emu_to_inch(shape.height),
            "right":  round((emu_to_inch(shape.left) or 0) + (emu_to_inch(shape.width) or 0), 3),
            "bottom": round((emu_to_inch(shape.top) or 0) + (emu_to_inch(shape.height) or 0), 3),
            "cx":     round((emu_to_inch(shape.left) or 0) + (emu_to_inch(shape.width) or 0) / 2, 3),
            "cy":     round((emu_to_inch(shape.top) or 0) + (emu_to_inch(shape.height) or 0) / 2, 3),
        },
    }

    # 占位符
    if shape.is_placeholder:
        ph = shape.placeholder_format
        info["placeholder"] = {"idx": ph.idx, "type": str(ph.type)}

    # 文本
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                continue
            runs_info = []
            for run in para.runs:
                rd = {"text": run.text}
                if run.font.size:
                    rd["font_size_pt"] = round(run.font.size / 12700, 1)
                if run.font.name:
                    rd["font_name"] = run.font.name
                if run.font.bold:
                    rd["bold"] = True
                c = get_color_str(run.font.color)
                if c:
                    rd["color"] = c
                runs_info.append(rd)
            paragraphs.append({
                "text":  para.text,
                "align": str(para.alignment) if para.alignment else None,
                "runs":  runs_info,
            })
        info["text_frame"] = {
            "full_text":  shape.text_frame.text,
            "word_wrap":  shape.text_frame.word_wrap,
            "paragraphs": paragraphs,
        }
        # 溢出风险
        info["overflow_risk"] = assess_overflow_risk(shape)

    # 表格
    if shape.has_table:
        table = shape.table
        rows_data = []
        for r_idx, row in enumerate(table.rows):
            rows_data.append([
                {"row": r_idx, "col": c_idx, "text": cell.text}
                for c_idx, cell in enumerate(row.cells)
            ])
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "data": rows_data,
        }

    # 图表
    if shape.has_chart:
        chart = shape.chart
        info["chart"] = {
            "chart_type":   str(chart.chart_type),
            "series_count": len(chart.series),
            "has_title":    chart.has_title,
            "title": chart.chart_title.text_frame.text if chart.has_title else None,
        }

    # 填充色
    try:
        c = get_color_str(shape.fill.fore_color)
        if c:
            info["fill_color"] = c
    except Exception:
        pass

    # Group shape：递归列出子shape名称
    try:
        if shape.shape_type == 7:
            child_names = []
            for child in shape.shapes:
                try:
                    child_names.append(child.name)
                except Exception:
                    pass
            info["group_children"] = child_names
    except Exception:
        pass

    return info


def analyze_relations(shapes_info, slide_w, slide_h):
    """
    空间邻近关联分析：
    - 识别固定边界元素（贯穿型边条、角落小元素）
    - 对每个 shape 找出最近的其他 shape（上/下/左/右）
    - 识别可能的分组（中心轴对齐的 shape 组）
    返回 dict
    """
    result = {
        "boundary_elements": [],
        "anchor_groups":     [],
        "neighbors":         [],
    }

    # 1. 识别固定边界元素
    for s in shapes_info:
        p = s["position"]
        h = p["height"] or 0
        w = p["width"] or 0
        l = p["left"] or 0
        t = p["top"] or 0

        reasons = []
        # 贯穿型：高度 >= 幻灯片高度80%
        if h >= slide_h * 0.8:
            reasons.append(f"贯穿型（高={h}\"，幻灯片高={slide_h}\"）")
        # 角落小元素：面积小且靠近角落
        area = w * h
        corner_margin = 1.0
        in_corner = (
            (l < corner_margin or l + w > slide_w - corner_margin) and
            (t < corner_margin or t + h > slide_h - corner_margin)
        )
        if area < 0.5 and in_corner and not s.get("text_frame"):
            reasons.append(f"角落装饰（面积={round(area,3)}\"²）")

        if reasons:
            result["boundary_elements"].append({
                "shape_id": s["shape_id"],
                "name":     s["name"],
                "reason":   " | ".join(reasons),
                "right":    p["right"],
                "bottom":   p["bottom"],
                "left":     l,
                "top":      t,
            })

    # 2. 识别竖向对齐组（cx 相近的 shape，可能是卡片+圆圈关联）
    cx_groups = {}
    for s in shapes_info:
        cx = s["position"]["cx"]
        matched = False
        for key in list(cx_groups.keys()):
            if abs(cx - key) < 0.3:  # 0.3in 容差
                cx_groups[key].append(s)
                matched = True
                break
        if not matched:
            cx_groups[cx] = [s]

    for cx, group in cx_groups.items():
        if len(group) >= 2:
            result["anchor_groups"].append({
                "center_x":  round(cx, 3),
                "shape_ids": [s["shape_id"] for s in group],
                "names":     [s["name"] for s in group],
            })

    # 3. 每个 shape 的最近邻（上下左右各一个）
    for s in shapes_info:
        sp = s["position"]
        scx, scy = sp["cx"], sp["cy"]
        neighbors = {"up": None, "down": None, "left": None, "right": None}
        min_dist  = {"up": 999, "down": 999, "left": 999, "right": 999}

        for other in shapes_info:
            if other["shape_id"] == s["shape_id"]:
                continue
            op  = other["position"]
            ocx = op["cx"]
            ocy = op["cy"]
            dx  = ocx - scx
            dy  = ocy - scy
            dist = math.sqrt(dx*dx + dy*dy)

            # 方向判断：以 45° 分区
            if abs(dy) > abs(dx):
                direction = "down" if dy > 0 else "up"
            else:
                direction = "right" if dx > 0 else "left"

            if dist < min_dist[direction]:
                min_dist[direction] = dist
                neighbors[direction] = {
                    "shape_id": other["shape_id"],
                    "name":     other["name"],
                    "dist":     round(dist, 3),
                }

        result["neighbors"].append({
            "shape_id":  s["shape_id"],
            "name":      s["name"],
            "neighbors": neighbors,
        })

    return result


def inspect_slide(slide, slide_num, slide_w_inch, slide_h_inch):
    result = {
        "slide_number": slide_num,
        "shape_count":  len(slide.shapes),
        "shapes":       [],
        "notes":        None,
        "background":   None,
    }
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            result["notes"] = notes_text
    try:
        c = get_color_str(slide.background.fill.fore_color)
        if c:
            result["background"] = c
    except Exception:
        pass

    for idx in range(len(slide.shapes)):
        try:
            shape = slide.shapes[idx]
            result["shapes"].append(inspect_shape(shape, idx))
        except Exception as e:
            result["shapes"].append({
                "index":    idx,
                "shape_id": None,
                "name":     f"[读取失败: {type(e).__name__}]",
                "type":     "UNKNOWN",
                "position": {"left": None, "top": None, "width": None, "height": None,
                             "right": None, "bottom": None, "cx": None, "cy": None},
                "error":    str(e),
            })

    # 关联分析附加到 slide 结果
    result["relations"] = analyze_relations(result["shapes"], slide_w_inch, slide_h_inch)
    return result


def inspect_layouts(prs):
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = [
            {"idx": ph.placeholder_format.idx,
             "type": str(ph.placeholder_format.type),
             "name": ph.name}
            for ph in layout.placeholders
        ]
        layouts.append({
            "index": i,
            "name":  layout.name,
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
        })
    return layouts


def print_report(report, show_relations=False, show_overflow=False):
    print(f"\n{'='*66}")
    print(f"📊 文件：{report['file']}")
    print(f"📄 总页数：{report['total_slides']}")
    print(f"📐 尺寸：{report['slide_width_inch']}\" × {report['slide_height_inch']}\"")

    if report.get("layouts"):
        print(f"\n📋 可用布局（共 {len(report['layouts'])} 个）：")
        for l in report["layouts"]:
            print(f"  [{l['index']}] {l['name']} — {l['placeholder_count']} 个占位符")

    for s in report["slides"]:
        print(f"\n{'─'*66}")
        print(f"📑 第 {s['slide_number']} 页  |  {s['shape_count']} 个 Shape", end="")
        if s.get("background"):
            print(f"  |  背景：{s['background']}", end="")
        print()
        if s.get("notes"):
            print(f"   📝 备注：{s['notes'][:80]}")

        for shape in s["shapes"]:
            p = shape["position"]
            ph_str   = f" [占位符 idx={shape['placeholder']['idx']}]" if shape.get("placeholder") else ""
            fill_str = f"  fill={shape['fill_color']}" if shape.get("fill_color") else ""
            print(f"\n  [{shape['index']}] id={shape.get('shape_id','?')} {shape['name']} | {shape['type']}{ph_str}")
            print(f"       📍 L={p['left']}\" T={p['top']}\"  W={p['width']}\" H={p['height']}\"  "
                  f"R={p['right']}\" B={p['bottom']}\"  cx={p['cx']}\" cy={p['cy']}\"{fill_str}")

            if shape.get("text_frame"):
                tf = shape["text_frame"]
                preview  = tf["full_text"][:80].replace('\n', ' / ')
                wrap_str = f"  wrap={'on' if tf.get('word_wrap') else 'off'}"
                print(f"       📝 文本：'{preview}'{wrap_str}")
                for para in tf["paragraphs"]:
                    for run in para.get("runs", []):
                        parts = []
                        if run.get("font_size_pt"): parts.append(f"{run['font_size_pt']}pt")
                        if run.get("font_name"):    parts.append(run["font_name"])
                        if run.get("bold"):         parts.append("bold")
                        if run.get("color"):        parts.append(run["color"])
                        if parts:
                            print(f"       🔤 字体：{' | '.join(parts)}")
                        break

                # 溢出风险：--overflow 时显示全部；不加时仅显示高风险警告
                if shape.get("overflow_risk"):
                    risk_info = shape["overflow_risk"]
                    if risk_info["risk"] == "high":
                        # 无论是否加 --overflow，高风险始终提示
                        print(f"       🔴 溢出风险：HIGH（建议检查）")
                        if show_overflow:
                            for d in risk_info["details"]:
                                print(f"          └ '{d['text']}' "
                                      f"估算宽={d['est_width']}\" 框宽={d['box_width']}\" "
                                      f"比率={d['ratio']}")
                    elif show_overflow and risk_info["risk"] == "medium":
                        print(f"       🟡 溢出风险：MEDIUM")
                        for d in risk_info["details"]:
                            print(f"          └ '{d['text']}' "
                                  f"估算宽={d['est_width']}\" 框宽={d['box_width']}\" "
                                  f"比率={d['ratio']}")

            if shape.get("table"):
                t = shape["table"]
                print(f"       📊 表格：{t['rows']}行 × {t['cols']}列")
                for row in t["data"][:3]:
                    print(f"         {[c['text'][:15] for c in row]}")
                if t["rows"] > 3:
                    print(f"         ... 共 {t['rows']} 行")

            if shape.get("chart"):
                c = shape["chart"]
                title_str = f"  标题：{c['title']}" if c.get("title") else ""
                print(f"       📈 图表：{c['chart_type']} | {c['series_count']} 系列{title_str}")

        # 关联分析输出
        if show_relations and s.get("relations"):
            rel = s["relations"]
            print(f"\n  {'─'*50}")
            print(f"  🔗 空间关联分析")

            if rel["boundary_elements"]:
                print(f"\n  【固定边界元素】")
                for b in rel["boundary_elements"]:
                    print(f"    id={b['shape_id']} '{b['name']}'  → {b['reason']}")
                    print(f"       L={b['left']}\" R={b['right']}\" T={b['top']}\" B={b['bottom']}\"")

            if rel["anchor_groups"]:
                print(f"\n  【竖向对齐组（可能的锚点关联）】")
                for g in rel["anchor_groups"]:
                    print(f"    中心X≈{g['center_x']}\"  shape_ids={g['shape_ids']}")
                    print(f"       names={g['names']}")

    print(f"\n{'='*66}")
    print("✅ 分析完成")


def main():
    parser = argparse.ArgumentParser(description="PPTX 结构完整分析工具")
    parser.add_argument("file",               help="PPTX 文件路径")
    parser.add_argument("--slide",   type=int, default=None, help="只分析第N页（1-based）")
    parser.add_argument("--layouts", action="store_true",    help="额外列出所有可用布局")
    parser.add_argument("--json",    action="store_true",    help="输出 JSON 格式")
    parser.add_argument("--relations", action="store_true",  help="输出 shape 空间关联分析")
    parser.add_argument("--overflow",  action="store_true",  help="输出文字溢出风险评估")
    args = parser.parse_args()

    pptx_path = Path(args.file)
    if not pptx_path.exists():
        print(f"❌ 文件不存在：{pptx_path}")
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    slide_w = emu_to_inch(prs.slide_width)
    slide_h = emu_to_inch(prs.slide_height)

    report = {
        "file":              str(pptx_path),
        "total_slides":      len(prs.slides),
        "slide_width_inch":  slide_w,
        "slide_height_inch": slide_h,
        "slides": [],
    }

    if args.layouts:
        report["layouts"] = inspect_layouts(prs)

    if args.slide:
        if args.slide < 1 or args.slide > len(prs.slides):
            print(f"❌ 页码超出范围（共 {len(prs.slides)} 页）")
            sys.exit(1)
        report["slides"].append(
            inspect_slide(prs.slides[args.slide - 1], args.slide, slide_w, slide_h))
    else:
        for i, slide in enumerate(prs.slides):
            report["slides"].append(inspect_slide(slide, i + 1, slide_w, slide_h))

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print_report(report,
                     show_relations=args.relations,
                     show_overflow=args.overflow)


if __name__ == "__main__":
    main()
